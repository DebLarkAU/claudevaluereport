#!/usr/bin/env python3
"""
Build a generic "TrendAI Vision One Value Report" deck for any customer.

Rerun workflow (see the customer's own runbook, generated from RUNBOOK_TEMPLATE.md):
  1. Refresh the figures in value_report_data.json (pull from the portal).
  2. python3 build_value_report.py
  -> writes  ~/Desktop/<customer> Value Report.pptx  (customer name comes from
     the "customer" field in value_report_data.json — no per-customer edits
     to this script are needed)

Design is fixed; only value_report_data.json changes between customers/runs.
Starts from template.pptx (title slide + all slide layouts, incl. "Quote 01").

This file is copied verbatim into each customer's own
~/Documents/<Customer>/value_report/ directory by the generic "value-report"
skill. If a customer's Vision One tenant has quirks that need bespoke handling
(an empty widget, an alternate dominant metric, etc. — see how Intralot's copy
diverged from John H Carter's over time), extend THAT customer's copy; don't
assume every tenant behaves the same way.
"""
import json, os, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "template.pptx")
DATA_FILE = os.path.join(HERE, "value_report_data.json")

# ---- palette ----
BG=RGBColor(0x14,0x18,0x1E); DARK=RGBColor(0x21,0x25,0x2B)
CARD=RGBColor(0x23,0x29,0x33); CARD2=RGBColor(0x1C,0x22,0x2B); CARDL=RGBColor(0xF4,0xF5,0xF7)
BORDER=RGBColor(0x3A,0x42,0x50); BORDERL=RGBColor(0xE3,0xE5,0xE9)
TXT=RGBColor(0xEC,0xEE,0xF1); INK=RGBColor(0x26,0x26,0x26); SUB=RGBColor(0x97,0x9E,0xAA)
GRAY=RGBColor(0x6E,0x6E,0x6E); LGRAY=RGBColor(0xC9,0xCE,0xD6); WHITE=RGBColor(0xFF,0xFF,0xFF)
RED=RGBColor(0xD0,0x00,0x1E); REDB=RGBColor(0xF0,0x55,0x5F); ORANGE=RGBColor(0xF2,0x99,0x4A)
ORANGE2=RGBColor(0xE8,0x8A,0x3A); TEAL=RGBColor(0x57,0xC4,0xD6); BLUE=RGBColor(0x4F,0x9F,0xE0)
GREEN=RGBColor(0x35,0xC7,0x59); RANKC=RGBColor(0xB6,0xBB,0xC4); TRACK=RGBColor(0x2C,0x33,0x3E)
GRID=RGBColor(0x33,0x3B,0x47); ROWLN=RGBColor(0x33,0x3B,0x47); PH=RGBColor(0x6E,0x77,0x86)

def rrect(s,x,y,w,h,fill,line=None,radius=0.08,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid();sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line;sp.line.width=Pt(0.75)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=radius
    except: pass
    return sp
def oval(s,x,y,w,h,fill,line=None,lw=1.5):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid();sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line;sp.line.width=Pt(lw)
    sp.shadow.inherit=False;return sp
def tbox(s,x,y,w,h,anchor=MSO_ANCHOR.MIDDLE):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame;tf.word_wrap=True;tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.04);tf.margin_right=Inches(0.04);tf.margin_top=Inches(0.01);tf.margin_bottom=Inches(0.01)
    return tf
def line(tf,runs,align=PP_ALIGN.LEFT,sb=0,first=False):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align;p.space_after=Pt(0)
    if sb:p.space_before=Pt(sb)
    for (t,sz,b,c) in runs:
        r=p.add_run();r.text=t;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name="Arial"
    return p
def clear(s):
    for sh in list(s.shapes): sh._element.getparent().remove(sh._element)
def header(s,title,subtitle,bg=True):
    if bg: rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[(title,30,True,WHITE)],first=True)
    rrect(s,0.55,0.93,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(subtitle,11,False,SUB)],first=True)
def headerbar(s,x,y,w,title,accent=DARK,h=0.5,fs=12.5):
    bar=rrect(s,x,y,w,h,accent,radius=0.05)
    tf=bar.text_frame;tf.vertical_anchor=MSO_ANCHOR.MIDDLE;tf.margin_left=Inches(0.16)
    p=tf.paragraphs[0];r=p.add_run();r.text=title;r.font.size=Pt(fs);r.font.bold=True;r.font.color.rgb=WHITE;r.font.name="Arial"

CUST=None  # set in main

# ========================= TITLE SLIDE =========================
def build_title(slide,D):
    def settext(name,text):
        for sh in slide.shapes:
            if sh.name==name and sh.has_text_frame:
                tf=sh.text_frame; p=tf.paragraphs[0]
                if p.runs:
                    p.runs[0].text=text
                    for ex in p.runs[1:]: ex._r.getparent().remove(ex._r)
                    run=p.runs[0]
                else:
                    run=p.add_run(); run.text=text
                run.font.color.rgb=WHITE
                for xp in tf.paragraphs[1:]: xp._p.getparent().remove(xp._p)
                return
    settext('Title 1', D['customer'])
    settext('Subtitle 2', D['report_title'])
    settext('Text Placeholder 3', "Presented by:  "+D['presenter'])
    # template.pptx may carry a stale leftover date textbox (e.g. "TextBox 4" /
    # "June 2026") baked in from whatever run it was last saved from. Clear any
    # text shape not one of the three known placeholders before adding the
    # fresh date box, or the old text overlaps/interleaves with the new text.
    known={'Title 1','Subtitle 2','Text Placeholder 3'}
    for sh in list(slide.shapes):
        if sh.has_text_frame and sh.name not in known:
            sh._element.getparent().remove(sh._element)
    tf=tbox(slide,6.67,4.45,6.04,0.45,MSO_ANCHOR.MIDDLE); tf.margin_left=Inches(0.02)
    line(tf,[(D['month_year'],18,False,WHITE)],first=True)
    slide.shapes[-1].name='TITLE_DATE'   # named so preserve-mode reruns can refresh just the date

def update_title_date(slide,D):
    """Preserve-mode helper: refresh ONLY the date on the (hand-edited) title slide,
    leaving the presenter name and everything else intact."""
    target=None
    for sh in slide.shapes:
        if sh.name=='TITLE_DATE' and sh.has_text_frame: target=sh; break
    if target is None:   # older decks: date box wasn't named yet — match by position
        for sh in slide.shapes:
            if sh.has_text_frame and sh.left is not None and abs(sh.left-Inches(6.67))<Inches(0.3) and abs(sh.top-Inches(4.45))<Inches(0.3):
                target=sh; target.name='TITLE_DATE'; break
    if target is None:   # date box was removed entirely — add a fresh one
        tf=tbox(slide,6.67,4.45,6.04,0.45,MSO_ANCHOR.MIDDLE); tf.margin_left=Inches(0.02)
        line(tf,[(D['month_year'],18,False,WHITE)],first=True)
        slide.shapes[-1].name='TITLE_DATE'; return
    tf=target.text_frame; p=tf.paragraphs[0]
    if p.runs:
        p.runs[0].text=D['month_year']
        for ex in p.runs[1:]: ex._r.getparent().remove(ex._r)
        r=p.runs[0]; r.font.size=Pt(18); r.font.color.rgb=WHITE; r.font.name="Arial"
    else:
        line(tf,[(D['month_year'],18,False,WHITE)],first=True)
    for xp in tf.paragraphs[1:]: xp._p.getparent().remove(xp._p)

# ========================= 1. CYBER RISK =========================
def build_cyber(s,D):
    d=D['cyber']
    # NOTE: relies on "Quote 01" dark layout background (no full-bleed rect)
    line(tbox(s,0.5,0.30,12.33,0.62,MSO_ANCHOR.TOP),[("Cyber Risk Exposure Management",30,True,WHITE)],first=True)
    rrect(s,0.55,0.95,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.02,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™   •   as of {D['as_of']}",12,False,LGRAY)],first=True)
    hero=rrect(s,0.5,1.48,12.33,1.45,DARK,radius=0.06)
    line(tbox(s,0.85,1.62,4.5,0.4),[("CYBER RISK INDEX",14,True,LGRAY)],first=True)
    line(tbox(s,0.80,1.92,3.4,0.95),[(d['cri'],54,True,ORANGE),("  /100",20,True,LGRAY)],first=True)
    chip=rrect(s,3.95,2.18,1.9,0.5,ORANGE,radius=0.5)
    ct=chip.text_frame;ct.vertical_anchor=MSO_ANCHOR.MIDDLE;cp=ct.paragraphs[0];cp.alignment=PP_ALIGN.CENTER
    cr=cp.add_run();cr.text=d['cri_level'];cr.font.size=Pt(13);cr.font.bold=True;cr.font.color.rgb=DARK;cr.font.name="Arial"
    tf=tbox(s,6.4,1.70,6.2,1.1)
    line(tf,[("Overall exposure across devices, accounts, vulnerabilities,",12.5,False,LGRAY)],first=True)
    line(tf,[(f"applications and cloud assets.  Regional average: {d['regional_avg']}",12.5,False,LGRAY)])
    # ---- 2x3 grid: one panel per risk category ----
    RISK_CHIP={"Low":(GREEN,WHITE),"Medium":(ORANGE,DARK),"High":(RED,WHITE)}
    def cat_panel(x,y,w,h,cat):
        hh=0.46
        headerbar(s,x,y,w,cat['title'],DARK,hh,11.5)
        lvl=cat.get('risk_level')
        if lvl:
            col,tcol=RISK_CHIP.get(lvl,(GRAY,WHITE))
            cw,chh=0.85,0.28; cx=x+w-cw-0.10; cyc=y+(hh-chh)/2
            chip=rrect(s,cx,cyc,cw,chh,col,radius=0.5)
            ctf=chip.text_frame;ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
            cpp=ctf.paragraphs[0];crr=cpp.add_run();crr.text=lvl;crr.font.size=Pt(8.5);crr.font.bold=True;crr.font.color.rgb=tcol;crr.font.name="Arial"
        cy2=y+hh+0.10; ch2=h-hh-0.10
        rrect(s,x,cy2,w,ch2,CARDL,BORDERL,radius=0.08)
        tf=tbox(s,x+0.1,cy2+0.06,w-0.2,ch2-0.12)
        line(tf,[(cat['headline'],26,True,RED)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(cat['sub'],10,False,GRAY)],align=PP_ALIGN.CENTER,sb=4)
    panel_y=3.05; col_gap=0.28; row_gap=0.20
    pw=(12.33-2*col_gap)/3.0; row_h=(7.35-panel_y-row_gap)/2.0
    xs=[0.5,0.5+pw+col_gap,0.5+2*(pw+col_gap)]
    ys=[panel_y,panel_y+row_h+row_gap]
    positions=[(xs[0],ys[0]),(xs[1],ys[0]),(xs[2],ys[0]),(xs[0],ys[1]),(xs[1],ys[1]),(xs[2],ys[1])]
    for (x,y),cat in zip(positions,d['categories']):
        cat_panel(x,y,pw,row_h,cat)

# ========================= 1b. RISK FACTORS DETAIL (1-2 per slide) =========================
def build_risk_detail_slide(s,D,cats):
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Risk Factors Detail",30,True,WHITE)],first=True)
    rrect(s,0.55,0.95,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    subtitle=" & ".join(c['title'] for c in cats)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™   •   {subtitle}   •   as of {D['as_of']}",11,False,LGRAY)],first=True)
    COL_Y=1.55; COL_H=5.45; gap=0.35
    n=len(cats); w=(12.33-(n-1)*gap)/n
    RISK_CHIP={"Low":(GREEN,WHITE),"Medium":(ORANGE,DARK),"High":(RED,WHITE)}
    for i,cat in enumerate(cats):
        x=0.5+i*(w+gap)
        rrect(s,x,COL_Y,w,COL_H,CARD,BORDER,radius=0.05)
        hh=0.52
        headerbar(s,x,COL_Y,w,cat['title'],DARK,hh,14)
        lvl=cat.get('risk_level')
        if lvl:
            col,tcol=RISK_CHIP.get(lvl,(GRAY,WHITE))
            cw,chh=1.0,0.32; cx=x+w-cw-0.14; cyc=COL_Y+(hh-chh)/2
            chip=rrect(s,cx,cyc,cw,chh,col,radius=0.5)
            ctf=chip.text_frame;ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
            cpp=ctf.paragraphs[0];crr=cpp.add_run();crr.text=lvl;crr.font.size=Pt(10);crr.font.bold=True;crr.font.color.rgb=tcol;crr.font.name="Arial"
        factors=cat.get('factors') or []
        if not factors:
            note=cat.get('note','No risk factors detected.')
            tf=tbox(s,x+0.3,COL_Y+hh+0.3,w-0.6,COL_H-hh-0.6,MSO_ANCHOR.MIDDLE)
            line(tf,[("✓  ",16,True,GREEN),(note,13,False,TXT)],first=True)
            continue
        row_y0=COL_Y+hh+0.10; row_bot=COL_Y+COL_H-0.10
        rh=(row_bot-row_y0)/len(factors)
        for j,(count,label) in enumerate(factors):
            ry=row_y0+j*rh
            if j%2==1: rrect(s,x+0.08,ry,w-0.16,rh,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
            line(tbox(s,x+0.22,ry,w-1.5-0.3,rh,MSO_ANCHOR.MIDDLE),[(label,12.5,False,TXT)],first=True)
            line(tbox(s,x+w-1.5-0.05,ry,1.5,rh,MSO_ANCHOR.MIDDLE),[(count,16,True,RED)],align=PP_ALIGN.RIGHT,first=True)
    line(tbox(s,0.5,7.08,12.33,0.3,MSO_ANCHOR.TOP),[("Risk factors are the underlying signals driving each category's risk level shown on the Cyber Risk Exposure Management overview.",8.5,False,SUB)],first=True)

# ========================= 2. MDR =========================
def build_mdr(s,D):
    d=D['mdr']; LG=LGRAY
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.32,12.33,0.62,MSO_ANCHOR.TOP),[("Managed Detection and Response",30,True,WHITE)],first=True)
    rrect(s,0.55,0.97,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    if not d.get('applicable',True):
        line(tbox(s,0.5,1.04,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ MDR",12,False,GRAY)],first=True)
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"Managed Detection and Response is not part of {CUST}'s TrendAI Vision One subscription.")
        line(tf,[(note,15,False,TXT)],align=PP_ALIGN.CENTER,first=True)
        return
    line(tbox(s,0.5,1.04,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ MDR   •   Reporting period: {d['period']}",12,False,GRAY)],first=True)
    # ---- left: vertical tapering funnel ----
    line(tbox(s,0.85,1.58,5.4,0.30,MSO_ANCHOR.TOP),[("MDR SECURITY EVENT FUNNEL",13,True,LG)],first=True)
    funnel=d['funnel']  # list of [num,label]; last one green
    blues=[RGBColor(0x1B,0x4B,0x7A),RGBColor(0x23,0x5E,0x9C),RGBColor(0x2E,0x77,0xBE),RGBColor(0x4F,0x9F,0xE0),RGBColor(0x7E,0xBB,0xEC)]
    cx=3.30; wtop=5.10; wbot=2.00; n=len(funnel); fh=0.62; gap=0.205; fy=2.12
    for i,(num,lab) in enumerate(funnel):
        last=(i==n-1)
        w=wtop-(wtop-wbot)*i/(n-1); x=cx-w/2.0; y=fy+i*(fh+gap)
        rrect(s,x,y,w,fh,GREEN if last else blues[min(i,len(blues)-1)],radius=0.14)
        col=WHITE if i<3 else INK
        line(tbox(s,x,y,w,fh,MSO_ANCHOR.MIDDLE),[(num,15,True,col),("   "+lab,10.5,False,col)],align=PP_ALIGN.CENTER,first=True)

    # ---- right column ----
    rx=6.60; rw=12.83-rx
    def dcard(x,y,w,h): rrect(s,x,y,w,h,CARD,BORDER,radius=0.10)
    # Response Time (MTTx)
    line(tbox(s,rx,1.58,rw,0.30,MSO_ANCHOR.TOP),[("RESPONSE TIME (MTTx)",13,True,LG)],first=True)
    cg=0.18; cw=(rw-2*cg)/3.0; cyt=2.04; cht=1.28; xs2=[rx,rx+cw+cg,rx+2*(cw+cg)]
    for (val,lab),x in zip([(d['mtta'],"Acknowledge"),(d['mtti'],"Investigate"),(d['mttr'],"Respond")],xs2):
        dcard(x,cyt,cw,cht)
        tf=tbox(s,x+0.04,cyt+0.05,cw-0.08,cht-0.10,MSO_ANCHOR.MIDDLE)
        line(tf,[(val,22,True,WHITE)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(lab,11.5,False,SUB)],align=PP_ALIGN.CENTER,sb=4)
    line(tbox(s,rx,cyt+cht+0.07,rw,0.26,MSO_ANCHOR.TOP),[(d.get('mttx_caption',""),9,False,GRAY)],first=True)
    # Response Action Taken
    line(tbox(s,rx,3.74,rw,0.30,MSO_ANCHOR.TOP),[("RESPONSE ACTION TAKEN",13,True,LG)],first=True)
    ay=4.14; ah=1.12; dcard(rx,ay,rw,ah)
    line(tbox(s,rx+0.26,ay+0.16,rw-1.7,0.36,MSO_ANCHOR.TOP),[(d.get('action_title',""),15,True,WHITE)],first=True)
    pillw=1.18; pillh=0.36; px=rx+rw-0.26-pillw; py=ay+0.18
    rrect(s,px,py,pillw,pillh,GREEN,radius=0.5)
    line(tbox(s,px,py,pillw,pillh,MSO_ANCHOR.MIDDLE),[(d.get('action_status',"APPROVED"),10.5,True,RGBColor(0x0E,0x3D,0x1E))],align=PP_ALIGN.CENTER,first=True)
    line(tbox(s,rx+0.26,ay+0.58,rw-0.5,0.46,MSO_ANCHOR.TOP),[(d.get('action_sub',""),10,False,SUB)],first=True)
    # Outcome
    line(tbox(s,rx,5.42,rw,0.30,MSO_ANCHOR.TOP),[("OUTCOME",13,True,LG)],first=True)
    oy=5.82; oh=1.05; dcard(rx,oy,rw,oh)
    line(tbox(s,rx+0.24,oy,0.72,oh,MSO_ANCHOR.MIDDLE),[(d.get('outcome_num',"0"),34,True,GREEN)],align=PP_ALIGN.LEFT,first=True)
    line(tbox(s,rx+0.98,oy+0.18,rw-1.25,0.36,MSO_ANCHOR.TOP),[(d.get('outcome_label',""),15,True,WHITE)],first=True)
    line(tbox(s,rx+0.98,oy+0.56,rw-1.25,0.36,MSO_ANCHOR.TOP),[(d.get('outcome_sub',""),10,False,SUB)],first=True)

# ========================= 2b. DATA SOURCE AND LOG MANAGEMENT =========================
def build_data_source(s,D):
    d=D['data_source']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("Data Source and Log Management",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   Agentic SIEM and XDR · Data Monitoring   •   {d.get('period','')}   •   as of {D['as_of']}",11,False,SUB)],first=True)
    rrect(s,0.5,1.4,12.33,1.05,CARD,BORDER,radius=0.05)
    metrics=d['metrics']  # list of {label,value}
    n=len(metrics); seg=12.33/n
    for i,m in enumerate(metrics):
        cx=0.5+seg*i+seg/2
        tf=tbox(s,cx-seg/2+0.1,1.43,seg-0.2,0.99)
        line(tf,[(m['value'],24,True,WHITE)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(m['label'],10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    for i in range(1,n):
        vx=0.5+seg*i
        rrect(s,vx,1.56,0.014,0.73,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    ty,th=2.65,4.45
    rrect(s,0.5,ty,12.33,th,CARD,BORDER,radius=0.04)
    headerbar(s,0.5,ty,12.33,"INGESTION AND RETENTION BY DATA SOURCE",RED,0.46,12.5)
    cols=[("Data Source",4.55),("Category",2.15),("Ingestion (GB)",2.75),("Retention (GB)",2.78)]
    hx=0.68; hy=ty+0.46+0.08
    for lab,w in cols:
        line(tbox(s,hx,hy,w,0.26),[(lab,9.5,True,SUB)],first=True)
        hx+=w
    sources=d['sources']  # list of [name, category, ingestion, retention]
    row_y0=hy+0.30; row_bot=ty+th-0.10
    rh=(row_bot-row_y0)/max(1,len(sources))
    for i,(name,cat,ing,ret) in enumerate(sources):
        ry=row_y0+i*rh
        if i%2==1: rrect(s,0.58,ry,12.17,rh,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
        cx=0.68
        line(tbox(s,cx,ry,4.55-0.1,rh,MSO_ANCHOR.MIDDLE),[(name,10.5,False,TXT)],first=True); cx+=4.55
        line(tbox(s,cx,ry,2.15-0.1,rh,MSO_ANCHOR.MIDDLE),[(cat,10,False,SUB)],first=True); cx+=2.15
        line(tbox(s,cx,ry,2.75-0.3,rh,MSO_ANCHOR.MIDDLE),[(ing,11,True,BLUE)],first=True); cx+=2.75
        line(tbox(s,cx,ry,2.78-0.3,rh,MSO_ANCHOR.MIDDLE),[(ret,11,True,ORANGE)],first=True)
    line(tbox(s,0.5,7.20,12.33,0.24,MSO_ANCHOR.TOP),[("Ingestion = new data received in the selected period. Retention = data retained (free + extended). Source: Agentic SIEM and XDR — Data Source and Log Management, Data Monitoring tab.",8,False,SUB)],first=True)

# ========================= 3. IDENTITY =========================
def build_identity(s,D):
    d=D['identity']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Identity Security Posture",30,True,WHITE)],first=True)
    rrect(s,0.55,0.95,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™   •   Risk level: {d['risk_level']}   •   Sources: AD + Microsoft Entra ID   •   as of {D['as_of']}",11,False,LGRAY)],first=True)
    COL_Y=1.52; COL_H=5.50; HDR_H=0.52; ROW_Y0=COL_Y+HDR_H+0.06; ROW_BOT=COL_Y+COL_H-0.06
    AMB=RGBColor(0xC2,0x73,0x00); TEAL_H=RGBColor(0x1F,0x6E,0x7E)
    AMB_V=RGBColor(0xE6,0xA8,0x3A); TEAL_V=TEAL; RED_V=REDB
    ATK_H=RGBColor(0x1E,0x4F,0x7A); ATK_V=BLUE
    def column(x,w,title,accent):
        rrect(s,x,COL_Y,w,COL_H,CARD,BORDER,radius=0.04)
        headerbar(s,x,COL_Y,w,title,accent,HDR_H,11)
    def rows(x,w,items,n,render):
        rh=(ROW_BOT-ROW_Y0)/n
        for i,it in enumerate(items):
            y=ROW_Y0+i*rh
            if i%2==1: rrect(s,x+0.06,y,w-0.12,rh,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
            render(x,w,y,rh,it)
    xA,wA=0.5,2.93; xB,wB=3.63,2.93; xD,wD=6.76,2.93; xC,wC=9.89,2.93
    column(xA,wA,"RISK EVENTS  ·  TOP 10",RED)
    def rA(x,w,y,rh,it):
        rk,name,imp,hc=it
        line(tbox(s,x+0.12,y,w-0.95-0.16,rh),[(rk+". ",9,True,RANKC),(name,9,False,TXT)],first=True)
        line(tbox(s,x+w-0.95-0.05,y,0.95,rh),[(imp,11,True,RED_V),(" ("+hc+")",7,False,SUB)],align=PP_ALIGN.RIGHT,first=True)
    rows(xA,wA,d['risk_events'],10,rA)
    column(xB,wB,"EXPOSURE EVENTS  ·  ALL",AMB)
    def rB(x,w,y,rh,it):
        name,imp,hc=it
        line(tbox(s,x+0.12,y,w-0.9-0.16,rh),[("• ",8.5,True,AMB_V),(name,8.5,False,TXT)],first=True)
        runs=[(imp,9.5,True,AMB_V)]
        if hc: runs.append((" ("+hc+")",7,False,SUB))
        line(tbox(s,x+w-0.9-0.05,y,0.9,rh),runs,align=PP_ALIGN.RIGHT,first=True)
    rows(xB,wB,d['exposure_events'],len(d['exposure_events']),rB)
    column(xD,wD,"ATTACK EVENTS",ATK_H)
    def rD(x,w,y,rh,it):
        name,imp,hc=it
        line(tbox(s,x+0.12,y,w-0.8-0.16,rh),[("• ",8.5,True,ATK_V),(name,8.5,False,TXT)],first=True)
        runs=[(imp,9.5,True,ATK_V)]
        if hc: runs.append((" ("+hc+")",7,False,SUB))
        line(tbox(s,x+w-0.8-0.05,y,0.8,rh),runs,align=PP_ALIGN.RIGHT,first=True)
    rows(xD,wD,d.get('attack_events',[]),max(1,len(d.get('attack_events',[]))),rD)
    column(xC,wC,"RISKY ACCOUNTS  ·  TOP 5",TEAL_H)
    risky=d.get('risky_accounts') or []
    if not risky:
        note=d.get('risky_accounts_note',"Risky-account ranking not available for this data source.")
        tf=tbox(s,xC+0.18,ROW_Y0,wC-0.36,ROW_BOT-ROW_Y0,MSO_ANCHOR.MIDDLE)
        line(tf,[(note,9.5,False,SUB)],align=PP_ALIGN.CENTER,first=True)
    else:
        def rC(x,w,y,rh,it):
            rk,name,sub,score=it
            ntf=tbox(s,x+0.12,y,w-0.8-0.16,rh)
            line(ntf,[(rk+". ",10,True,RANKC),(name,10.5,True,TXT)],first=True)
            line(ntf,[(sub,8,False,SUB)],sb=1)
            line(tbox(s,x+w-0.8-0.05,y,0.8,rh),[(score,12.5,True,TEAL_V)],align=PP_ALIGN.RIGHT,first=True)
        rows(xC,wC,risky,5,rC)
    line(tbox(s,0.5,7.12,12.33,0.3,MSO_ANCHOR.TOP),[("Impacted assets shown; (n) = highly-critical assets.  Attack & exposure events from Identity Security Posture.  Risky accounts ranked by asset risk score.",8.5,False,SUB)],first=True)

# ========================= 3b. DATA SECURITY =========================
def build_data_security(s,D):
    d=D['data_security']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Data Security",30,True,WHITE)],first=True)
    rrect(s,0.55,0.95,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Data Security Posture   •   Risk level: {d['risk_level']} ({d['score']}/100)   •   as of {D['as_of']}",11,False,LGRAY)],first=True)
    rrect(s,0.5,1.42,12.33,1.0,CARD,BORDER,radius=0.05)
    def hstat(cx,num,lab,col):
        tf=tbox(s,cx-1.95,1.46,3.9,0.92); line(tf,[(num,26,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    hstat(2.55,d['total_assets'],"Total assets",WHITE)
    hstat(6.665,d['sensitive_assets'],"Assets with sensitive data",REDB)
    hstat(10.78,d['monitored_assets'],"Monitored assets",TEAL)
    for vx in [4.61,8.72]: rrect(s,vx,1.58,0.014,0.68,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    px,pw=0.5,5.9; py,ph=2.6,4.55
    rrect(s,px,py,pw,ph,CARD,BORDER,radius=0.04)
    headerbar(s,px,py,pw,"SENSITIVE DATA DETECTIONS BY CATEGORY",RED,0.5,11.5)
    cats=d['categories']  # [name,count]
    maxv=float(max([c for _,c in cats]+[1]))
    bx0=px+1.55; bmaxw=pw-1.55-0.75
    top=py+0.78; avail=ph-0.78-0.2; rh=avail/len(cats)
    for i,(name,val) in enumerate(cats):
        cy=top+i*rh
        line(tbox(s,px+0.14,cy,1.35,rh),[(name,10.5,False,TXT)],first=True)
        rrect(s,bx0,cy+rh/2-0.13,bmaxw,0.26,TRACK,radius=0.3)
        bw=max(0.0,(val/maxv)*bmaxw)
        if bw>0.04: rrect(s,bx0,cy+rh/2-0.13,bw,0.26,RED,radius=0.3)
        line(tbox(s,bx0+max(bw,0.0)+0.08,cy,0.65,rh),[(str(val),12,True,WHITE)],first=True)
    ex,ew=6.75,6.08
    rrect(s,ex,py,ew,ph,CARD,BORDER,radius=0.04)
    headerbar(s,ex,py,ew,"TOP RISKY ASSETS WITH SENSITIVE DATA",TEAL,0.5,11.5)
    cols=[("Asset Name",1.85),("Type",1.05),("Sensitive Data Type",2.18),("Risk Score",0.9)]
    hx=ex+0.18; hy=py+0.5+0.08
    for lab,w in cols:
        line(tbox(s,hx,hy,w,0.24),[(lab,9,True,SUB)],first=True)
        hx+=w
    assets=d['top_assets']  # [name,type,sensitive_types,score]
    row_y0=hy+0.28; row_bot=py+ph-0.14
    rh2=(row_bot-row_y0)/max(1,len(assets))
    for i,(name,atype,stypes,score) in enumerate(assets):
        ry=row_y0+i*rh2
        if i%2==1: rrect(s,ex+0.08,ry,ew-0.16,rh2,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
        cx=ex+0.18
        line(tbox(s,cx,ry,1.85-0.1,rh2,MSO_ANCHOR.MIDDLE),[(name,10,True,TXT)],first=True); cx+=1.85
        line(tbox(s,cx,ry,1.05-0.1,rh2,MSO_ANCHOR.MIDDLE),[(atype,9.5,False,SUB)],first=True); cx+=1.05
        line(tbox(s,cx,ry,2.18-0.1,rh2,MSO_ANCHOR.MIDDLE),[(stypes,9.5,False,TXT)],first=True); cx+=2.18
        line(tbox(s,cx,ry,0.9-0.1,rh2,MSO_ANCHOR.MIDDLE),[(score,11,True,REDB)],first=True)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: Cyber Risk Exposure Management — Data Security Posture. Categories and top assets reflect the current sensitive-data scan.",8,False,SUB)],first=True)

# ========================= 4a. STANDARD ENDPOINT PROTECTION =========================
def build_standard_endpoint_protection(s,D):
    d=D.get('standard_endpoint_protection',{})
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Standard Endpoint Protection",30,True,WHITE)],first=True)
    rrect(s,0.55,0.93,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Standard Endpoint Protection   •   as of {D['as_of']}",11,False,SUB)],first=True)
    if not d.get('applicable',False):
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"Standard Endpoint Protection is not currently in use by {CUST}.")
        line(tf,[("Not Applicable",20,True,SUB)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(note,13,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
        line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Standard Endpoint Protection (Protection Manager), Endpoint Inventory.",8,False,SUB)],first=True)
        return
    es=d.get('endpoint_status',{})
    rrect(s,0.5,1.42,12.33,1.0,CARD,BORDER,radius=0.05)
    def hstat(cx,num,lab,col):
        tf=tbox(s,cx-1.95,1.46,3.9,0.92); line(tf,[(num,26,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    hstat(2.04,es.get('managed','0'),"Managed endpoints",WHITE)
    hstat(5.15,es.get('at_risk','0'),"At risk",REDB)
    hstat(8.26,es.get('outdated','0'),"Pattern/component outdated",ORANGE2)
    hstat(11.37,es.get('offline','0'),"Offline",GRAY)
    for vx in [3.59,6.70,9.81]: rrect(s,vx,1.58,0.014,0.68,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    px,pw=0.5,11.33; py,ph=2.6,4.55
    rrect(s,px,py,pw,ph,CARD,BORDER,radius=0.04)
    headerbar(s,px,py,pw,"THREAT DETECTIONS  ·  LAST 30 DAYS",RED,0.5,12.5)
    threats=d.get('threats') or []
    if not threats:
        tf=tbox(s,px+0.3,py+0.6,pw-0.6,ph-0.9,MSO_ANCHOR.MIDDLE)
        line(tf,[("✓  ",16,True,GREEN),("No threats detected on managed endpoints in the last 30 days.",13,False,TXT)],first=True)
    else:
        maxv=float(max([c for _,c in threats]+[1])); bx0=px+2.0; bmaxw=pw-2.0-0.9
        top=py+0.72; avail=ph-0.72-0.18; rh=avail/len(threats)
        for i,(name,val) in enumerate(threats):
            cy=top+i*rh
            line(tbox(s,px+0.16,cy,1.75,rh),[(name,10.5,False,TXT)],first=True)
            rrect(s,bx0,cy+rh/2-0.13,bmaxw,0.26,TRACK,radius=0.3)
            bw=max(0.0,(val/maxv)*bmaxw)
            if bw>0.04: rrect(s,bx0,cy+rh/2-0.13,bw,0.26,RED,radius=0.3)
            line(tbox(s,bx0+max(bw,0.0)+0.08,cy,0.75,rh),[(str(val),12,True,WHITE)],first=True)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Standard Endpoint Protection (Protection Manager), Endpoint Inventory.",8,False,SUB)],first=True)

# ========================= 4b. SERVER & WORKLOAD PROTECTION =========================
def build_server_workload_protection(s,D):
    d=D.get('server_workload_protection',{})
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Server & Workload Protection",30,True,WHITE)],first=True)
    rrect(s,0.55,0.93,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Server & Workload Protection   •   as of {D['as_of']}",11,False,SUB)],first=True)
    if not d.get('applicable',False):
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"Server & Workload Protection is not currently in use by {CUST}.")
        line(tf,[("Not Applicable",20,True,SUB)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(note,13,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
        line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Server & Workload Protection.",8,False,SUB)],first=True)
        return
    cs=d.get('computer_status',{})
    rrect(s,0.5,1.42,12.33,1.0,CARD,BORDER,radius=0.05)
    def hstat(cx,num,lab,col):
        tf=tbox(s,cx-1.95,1.46,3.9,0.92); line(tf,[(num,26,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    hstat(2.04,cs.get('managed','0'),"Managed computers",WHITE)
    hstat(5.15,cs.get('critical','0'),"Critical",REDB)
    hstat(8.26,cs.get('warning','0'),"Warning",ORANGE2)
    hstat(11.37,cs.get('unmanaged','0'),"Unmanaged",GRAY)
    for vx in [3.59,6.70,9.81]: rrect(s,vx,1.58,0.014,0.68,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    px,pw=0.5,12.33; py,ph=2.6,4.55
    rrect(s,px,py,pw,ph,CARD,BORDER,radius=0.04)
    al=d.get('alert_status',{})
    headerbar(s,px,py,pw,f"ALERT STATUS  ·  SECURITY EVENTS  ·  {al.get('period','LAST 30 DAYS').upper()}",RED,0.5,12)
    cats=al.get('categories') or []
    if not cats:
        tf=tbox(s,px+0.3,py+0.6,pw-0.6,ph-0.9,MSO_ANCHOR.MIDDLE)
        line(tf,[("No security events recorded.",13,False,TXT)],align=PP_ALIGN.CENTER,first=True)
    else:
        cols=[("Category",3.4),("Events",1.6),("Action Taken",6.7)]
        hx=px+0.24; hy=py+0.58
        for lab,w in cols:
            line(tbox(s,hx,hy,w,0.24),[(lab,9,True,SUB)],first=True)
            hx+=w
        row_y0=hy+0.30; row_bot=py+ph-0.14; rh=(row_bot-row_y0)/len(cats)
        for i,(name,count,action) in enumerate(cats):
            ry=row_y0+i*rh
            if i%2==1: rrect(s,px+0.08,ry,pw-0.16,rh,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
            capped="+" in str(count)
            cx=px+0.24
            line(tbox(s,cx,ry,3.4-0.1,rh,MSO_ANCHOR.MIDDLE),[(name,11,False,TXT)],first=True); cx+=3.4
            line(tbox(s,cx,ry,1.6-0.3,rh,MSO_ANCHOR.MIDDLE),[(str(count),12,True,(REDB if capped else WHITE))],first=True); cx+=1.6
            line(tbox(s,cx,ry,6.7-0.3,rh,MSO_ANCHOR.MIDDLE),[(action,10.5,False,SUB)],first=True)
        if any("+" in str(c) for _,c,_ in cats):
            line(tbox(s,px+0.24,py+ph-0.32,pw-0.48,0.24),[("+ = API result cap reached; actual count is higher.",7.5,False,SUB)],first=True)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Endpoint Event Viewer, Security Events (Server & Workload Protection agents).",8,False,SUB)],first=True)

# ========================= 4c. CLOUD SECURITY =========================
def build_cloud_security(s,D):
    d=D['cloud_security']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Cloud Security",30,True,WHITE)],first=True)
    rrect(s,0.55,0.93,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.0,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Cloud Security Posture   •   Risk level: {d['risk_level']} ({d['risk_index']}/100)   •   as of {D['as_of']}",11,False,LGRAY)],first=True)
    panel_y=1.5; col_gap=0.3; row_gap=0.25
    pw=(12.33-col_gap)/2.0; row_h=(7.15-panel_y-row_gap)/2.0
    xs=[0.5,0.5+pw+col_gap]; ys=[panel_y,panel_y+row_h+row_gap]
    HDR_H=0.42
    def panel(x,y,title,accent,p):
        rrect(s,x,y,pw,row_h,CARD,BORDER,radius=0.05)
        headerbar(s,x,y,pw,title,accent,HDR_H,11.5)
        hy=y+HDR_H+0.10
        tf=tbox(s,x+0.22,hy,pw-0.44,0.62)
        line(tf,[(p['headline_value'],26,True,accent)],first=True)
        line(tbox(s,x+0.22+1.7,hy+0.06,pw-0.44-1.7,0.5,MSO_ANCHOR.MIDDLE),[(p['headline_label'],11,True,TXT)],first=True)
        line(tbox(s,x+0.22,hy+0.60,pw-0.44,0.28),[(p['headline_sub'],9,False,SUB)],first=True)
        metrics=p.get('metrics') or []
        row_y0=hy+0.98; row_bot=y+row_h-0.10; rh=(row_bot-row_y0)/max(1,len(metrics))
        for i,(lab,val) in enumerate(metrics):
            ry=row_y0+i*rh
            if i%2==1: rrect(s,x+0.1,ry,pw-0.2,rh,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
            line(tbox(s,x+0.24,ry,pw-3.2,rh,MSO_ANCHOR.MIDDLE),[(lab,10,False,SUB)],first=True)
            line(tbox(s,x+pw-2.9,ry,2.7,rh,MSO_ANCHOR.MIDDLE),[(val,10.5,True,TXT)],align=PP_ALIGN.RIGHT,first=True)
    panel(xs[0],ys[0],"CLOUD OVERVIEW",RED,d['cloud_overview'])
    panel(xs[1],ys[0],"ENTITLEMENTS",TEAL,d['entitlements'])
    panel(xs[0],ys[1],"AI - SECURITY POSTURE MANAGEMENT",BLUE,d['ai_spm'])
    panel(xs[1],ys[1],"APIs",ORANGE2,d['apis'])
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: Cyber Risk Exposure Management — Cloud Security Posture (Cloud Overview, Entitlements, AI - Security Posture Management, APIs tabs).",8,False,SUB)],first=True)

# ========================= 4d. NETWORK SECURITY =========================
def build_network_security(s,D):
    d=D.get('network_security',{})
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("Network Security",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Network Security   •   as of {D['as_of']}",11,False,SUB)],first=True)
    if not d.get('applicable',False):
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"Network Security is not currently in use by {CUST}.")
        line(tf,[("Not Applicable",20,True,SUB)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(note,13,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
        line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Network Security, Network Overview / Network Inventory.",8,False,SUB)],first=True)
        return
    modules=d.get('modules') or []
    COL_Y=1.55; COL_H=5.55
    n=max(1,len(modules)); gap=0.3; cw=(12.33-(n-1)*gap)/n
    STATUS_CHIP={"In Use":(GREEN,WHITE),"Not Configured":(GRAY,WHITE)}
    for i,mod in enumerate(modules):
        x=0.5+i*(cw+gap)
        rrect(s,x,COL_Y,cw,COL_H,CARD,BORDER,radius=0.05)
        hh=0.5
        headerbar(s,x,COL_Y,cw,mod.get('name',''),DARK,hh,12)
        status=mod.get('status','')
        col,tcol=STATUS_CHIP.get(status,(GRAY,WHITE))
        cw2,chh=1.3,0.3; cx=x+cw-cw2-0.14; cyc=COL_Y+(hh-chh)/2
        chip=rrect(s,cx,cyc,cw2,chh,col,radius=0.5)
        ctf=chip.text_frame;ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
        cpp=ctf.paragraphs[0];crr=cpp.add_run();crr.text=status;crr.font.size=Pt(9);crr.font.bold=True;crr.font.color.rgb=tcol;crr.font.name="Arial"
        tf=tbox(s,x+0.24,COL_Y+hh+0.2,cw-0.48,COL_H-hh-0.4)
        for j,(ml,mv) in enumerate(mod.get('metrics') or []):
            line(tf,[(mv,20,True,WHITE)],first=(j==0),sb=(0 if j==0 else 12))
            line(tf,[(ml,10,False,SUB)],sb=2)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Network Security, Network Overview / Network Inventory.",8,False,SUB)],first=True)

# ========================= 5. CLOUD EMAIL =========================
def build_email(s,D):
    d=D['email']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.6,MSO_ANCHOR.TOP),[("Cloud Email & Collaboration Protection",28,True,WHITE)],first=True)
    rrect(s,0.55,0.92,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.99,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Email and Collaboration Security — Overview   •   as of {D['as_of']}",11,False,SUB)],first=True)
    if not d.get('applicable',False):
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"No email security data is currently available for {CUST}.")
        line(tf,[("Not Applicable",20,True,SUB)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(note,13,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
        line(tbox(s,0.5,7.24,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Email and Collaboration Security, Configuration and Operations — Overview.",8,False,SUB)],first=True)
        return
    rrect(s,0.5,1.42,12.33,1.18,CARD,BORDER,radius=0.05)
    def hstat(cx,num,lab,col):
        tf=tbox(s,cx-2.0,1.46,4.0,1.10); line(tf,[(num,29,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,11,False,SUB)],align=PP_ALIGN.CENTER,sb=3)
    hstat(2.55,d['scanned'],"Messages & files scanned",WHITE)
    hstat(6.66,d['threats_total'],"Total threat detections",REDB)
    hstat(10.77,d['top_threat'],d.get('top_threat_label',"Phishing — top threat"),ORANGE2)
    rrect(s,4.61,1.66,0.014,0.70,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    rrect(s,8.72,1.66,0.014,0.70,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    lx,lw=0.5,7.45; py=2.78; ph=4.45
    rrect(s,lx,py,lw,ph,CARD,BORDER,radius=0.04)
    headerbar(s,lx,py,lw,"THREAT DETECTIONS  ·  LAST 30 DAYS",RED,0.5,12.5)
    threats=d['threats']  # [name,val]; sorted desc
    maxv=float(d['threat_max']); bx0=lx+1.85; bmaxw=lw-1.85-0.85
    top=py+0.72; avail=ph-0.72-0.18; rh=avail/len(threats)
    for i,(name,val) in enumerate(threats):
        cy=top+i*rh
        line(tbox(s,lx+0.14,cy,1.65,rh),[(name,10.5,False,TXT)],first=True)
        rrect(s,bx0,cy+rh/2-0.13,bmaxw,0.26,TRACK,radius=0.3)
        bw=max(0.0,(val/maxv)*bmaxw)
        if bw>0.04: rrect(s,bx0,cy+rh/2-0.13,bw,0.26,ORANGE2,radius=0.3)
        line(tbox(s,bx0+max(bw,0.0)+0.08,cy,0.75,rh),[(str(val),12,True,WHITE)],first=True)
    rx,rw=8.15,4.68
    rrect(s,rx,py,rw,ph,CARD,BORDER,radius=0.04)
    headerbar(s,rx,py,rw,"TOP 5 HIGH-RISK EMAIL RECIPIENTS",TEAL,0.5,12)
    recips=d['recipients']  # [rank,email,sub,total]
    if not recips:
        line(tbox(s,rx+0.22,py+0.82,rw-0.44,1.4,MSO_ANCHOR.TOP),[("No individual high-risk email recipients were identified in the last 30 days.",11,False,SUB)],first=True)
    else:
        ry0=py+0.66; rav=ph-0.66-0.16; rrh=rav/len(recips)
        for i,(rk,email,sub,tot) in enumerate(recips):
            cy=ry0+i*rrh
            if i>0: rrect(s,rx+0.14,cy,rw-0.28,0.012,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
            ntf=tbox(s,rx+0.16,cy,rw-1.05,rrh)
            line(ntf,[(rk+".  ",10.5,True,RANKC),(email,10.5,True,TXT)],first=True)
            line(ntf,[(sub,9,False,SUB)],sb=2)
            line(tbox(s,rx+rw-0.95,cy,0.8,rrh),[(tot,17,True,REDB)],align=PP_ALIGN.RIGHT,first=True)
    line(tbox(s,0.5,7.24,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Email and Collaboration Security, Configuration and Operations — Overview. Total count = high-risk emails received per user (30 days).",8,False,SUB)],first=True)

# ========================= 6. INTELLIGENCE REPORTS =========================
def build_intel(s,D):
    d=D['intel']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.26,12.33,0.58,MSO_ANCHOR.TOP),[("Intelligence Reports",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Threat Intelligence   •   Matched sweeps · last 30 days   •   as of {D['as_of']}",11,False,SUB)],first=True)
    rrect(s,0.5,1.38,12.33,1.0,CARD,BORDER,radius=0.05)
    def hstat(cx,num,lab,col):
        tf=tbox(s,cx-1.9,1.42,3.8,0.92); line(tf,[(num,27,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    sm=d['summary']
    hstat(2.04,sm[0],"Matched reports",WHITE); hstat(5.15,sm[1],"Matched indicators",REDB)
    hstat(8.26,sm[2],"Affected assets",TEAL); hstat(11.37,sm[3],"Alerts triggered",GREEN)
    for vx in [3.59,6.70,9.81]: rrect(s,vx,1.54,0.014,0.68,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    cards=d['reports']  # [name,sub,src,sweeps,inds,assets,date,colorkey]
    colmap={"blue":BLUE,"teal":TEAL}
    cy=2.5; chh=3.55; cw=5.96; gap=0.41; xs=[0.5,0.5+cw+gap]
    if not cards:
        rrect(s,0.5,cy,12.33,chh,CARD,BORDER,radius=0.04)
        tf=tbox(s,1.1,cy+0.35,11.13,chh-0.7,MSO_ANCHOR.MIDDLE)
        line(tf,[("✓  No threat-intelligence matches in the last 30 days",17,True,GREEN)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(d.get('note',''),11.5,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
    for (name,sub,src,sweeps,inds,assets,dt,ck),cx in zip(cards,xs):
        accent=colmap.get(ck,BLUE)
        rrect(s,cx,cy,cw,chh,CARD,BORDER,radius=0.04)
        headerbar(s,cx,cy,cw,"MATCHED INTELLIGENCE REPORT",accent,0.46,10.5)
        tf=tbox(s,cx+0.25,cy+0.54,cw-0.5,0.82,MSO_ANCHOR.TOP)
        line(tf,[(name,15.5,True,WHITE)],first=True); line(tf,[(sub,9.5,False,SUB)],sb=2)
        oval(s,cx+0.27,cy+1.45,0.12,0.12,GREEN)
        line(tbox(s,cx+0.45,cy+1.34,cw-0.7,0.36),[("Source:  ",9.5,False,SUB),(src,10.5,True,TXT)],first=True)
        mw=(cw-0.5-2*0.18)/3; my=cy+1.78; mh=0.92
        for i,(ml,mv) in enumerate([("Matched sweeps",sweeps),("Matched indicators",inds),("Associated assets",assets)]):
            mx=cx+0.25+i*(mw+0.18)
            rrect(s,mx,my,mw,mh,CARD2,BORDER,radius=0.1)
            mtf=tbox(s,mx+0.04,my+0.05,mw-0.08,mh-0.1)
            line(mtf,[(mv,20,True,accent)],align=PP_ALIGN.CENTER,first=True); line(mtf,[(ml,8,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
        line(tbox(s,cx+0.25,cy+2.82,cw-3.05,0.34),[("Matched:  ",9.5,False,SUB),(dt,10.5,True,TXT)],first=True)
        badge=rrect(s,cx+cw-2.75,cy+2.82,2.5,0.42,RGBColor(0x1E,0x3A,0x28),radius=0.5)
        bt=badge.text_frame;bt.vertical_anchor=MSO_ANCHOR.MIDDLE;bp=bt.paragraphs[0];bp.alignment=PP_ALIGN.CENTER
        brn=bp.add_run();brn.text="✓  No alert triggered";brn.font.size=Pt(10.5);brn.font.bold=True;brn.font.color.rgb=GREEN;brn.font.name="Arial"
    ey=6.28; eh=0.98
    rrect(s,0.5,ey,12.33,eh,CARD2,BORDER,radius=0.05)
    rrect(s,0.5,ey,0.07,eh,GREEN,radius=0,shape=MSO_SHAPE.RECTANGLE)
    itf=tbox(s,0.78,ey+0.06,12.0,eh-0.12)
    line(itf,[("Why no alert was triggered",11,True,GREEN)],first=True)
    line(itf,[(d['explanation'],9.5,False,TXT)],sb=2)

# ========================= 6b. WORKFLOW AND AUTOMATION =========================
def build_workflow_automation(s,D):
    d=D.get('workflow_automation',{})
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("Workflow and Automation",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Workflow and Automation   •   as of {D['as_of']}",11,False,SUB)],first=True)
    COL_Y=1.55; COL_H=5.55; gap=0.41
    cw=(12.33-gap)/2; xs=[0.5,0.5+cw+gap]; HDR_H=0.5
    def empty_note(x,note):
        tf=tbox(s,x+0.3,COL_Y+HDR_H,cw-0.6,COL_H-HDR_H-0.2,MSO_ANCHOR.MIDDLE)
        line(tf,[("✓  ",16,True,GREEN),(note,13,False,TXT)],first=True)
    def rowlist(x,items,accent):
        row_y0=COL_Y+HDR_H+0.10; row_bot=COL_Y+COL_H-0.10
        rh=(row_bot-row_y0)/len(items)
        for i,name in enumerate(items):
            ry=row_y0+i*rh
            if i%2==1: rrect(s,x+0.08,ry,cw-0.16,rh,RGBColor(0x2B,0x32,0x3D),radius=0,shape=MSO_SHAPE.RECTANGLE)
            line(tbox(s,x+0.24,ry,cw-0.48,rh,MSO_ANCHOR.MIDDLE),[("•  ",12,True,accent),(name,12,False,TXT)],first=True)
    # ---- Security Playbooks ----
    x=xs[0]
    rrect(s,x,COL_Y,cw,COL_H,CARD,BORDER,radius=0.05)
    headerbar(s,x,COL_Y,cw,"SECURITY PLAYBOOKS",RED,HDR_H,13)
    playbooks=d.get('playbooks') or []
    if not playbooks:
        empty_note(x,d.get('playbooks_note',f"No Security Playbooks are currently configured for {CUST}."))
    else:
        rowlist(x,playbooks,REDB)
    # ---- Third-Party Integrations ----
    x=xs[1]
    rrect(s,x,COL_Y,cw,COL_H,CARD,BORDER,radius=0.05)
    headerbar(s,x,COL_Y,cw,"THIRD-PARTY INTEGRATIONS  ·  CONFIGURED",TEAL,HDR_H,12.5)
    integrations=d.get('integrations') or []
    if not integrations:
        empty_note(x,d.get('integrations_note',f"No third-party integrations are currently configured for {CUST}."))
    else:
        rowlist(x,integrations,TEAL)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Security Playbooks: Workflow and Automation — Security Playbooks. Integrations: Cyber Risk Exposure Management — Data Sources (status: Configured).",8,False,SUB)],first=True)

# ========================= 6c. AI SECURITY =========================
def build_ai_security(s,D):
    d=D.get('ai_security',{})
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("AI Security",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ AI Security   •   as of {D['as_of']}",11,False,SUB)],first=True)
    if not d.get('applicable',False):
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"AI Security (AI Security Blueprint, AI Application Security, AI Secure Access) is not currently in use by {CUST}.")
        line(tf,[("Not Applicable",20,True,SUB)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(note,13,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
        line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — AI Security Blueprint, AI Application Security, AI Secure Access.",8,False,SUB)],first=True)
        return
    modules=d.get('modules') or []
    COL_Y=1.55; COL_H=5.55
    n=max(1,len(modules)); gap=0.3; cw=(12.33-(n-1)*gap)/n
    STATUS_CHIP={"In Use":(GREEN,WHITE),"Not Configured":(GRAY,WHITE)}
    for i,mod in enumerate(modules):
        x=0.5+i*(cw+gap)
        rrect(s,x,COL_Y,cw,COL_H,CARD,BORDER,radius=0.05)
        hh=0.5
        headerbar(s,x,COL_Y,cw,mod.get('name',''),DARK,hh,12)
        status=mod.get('status','')
        col,tcol=STATUS_CHIP.get(status,(GRAY,WHITE))
        cw2,chh=1.3,0.3; cx=x+cw-cw2-0.14; cyc=COL_Y+(hh-chh)/2
        chip=rrect(s,cx,cyc,cw2,chh,col,radius=0.5)
        ctf=chip.text_frame;ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
        cpp=ctf.paragraphs[0];crr=cpp.add_run();crr.text=status;crr.font.size=Pt(9);crr.font.bold=True;crr.font.color.rgb=tcol;crr.font.name="Arial"
        tf=tbox(s,x+0.24,COL_Y+hh+0.2,cw-0.48,COL_H-hh-0.4)
        for j,(ml,mv) in enumerate(mod.get('metrics') or []):
            line(tf,[(mv,20,True,WHITE)],first=(j==0),sb=(0 if j==0 else 12))
            line(tf,[(ml,10,False,SUB)],sb=2)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — AI Security Blueprint, AI Application Security, AI Secure Access.",8,False,SUB)],first=True)

# ========================= 6d. ZERO TRUST =========================
def build_zero_trust(s,D):
    d=D.get('zero_trust',{})
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("Zero Trust",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI Vision One™ Zero Trust Secure Access   •   as of {D['as_of']}",11,False,SUB)],first=True)
    if not d.get('applicable',False):
        rrect(s,0.5,2.6,12.33,2.3,CARD,BORDER,radius=0.08)
        tf=tbox(s,0.9,2.6,11.5,2.3,MSO_ANCHOR.MIDDLE)
        note=d.get('not_applicable_note',f"Zero Trust Secure Access (Internet Access, Private Access, AI Secure Access) is not currently in use by {CUST}.")
        line(tf,[("Not Applicable",20,True,SUB)],align=PP_ALIGN.CENTER,first=True)
        line(tf,[(note,13,False,TXT)],align=PP_ALIGN.CENTER,sb=10)
        line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Zero Trust Secure Access, Secure Access Overview.",8,False,SUB)],first=True)
        return
    modules=d.get('modules') or []
    COL_Y=1.55; COL_H=5.55
    n=max(1,len(modules)); gap=0.3; cw=(12.33-(n-1)*gap)/n
    STATUS_CHIP={"In Use":(GREEN,WHITE),"Not Configured":(GRAY,WHITE)}
    for i,mod in enumerate(modules):
        x=0.5+i*(cw+gap)
        rrect(s,x,COL_Y,cw,COL_H,CARD,BORDER,radius=0.05)
        hh=0.5
        headerbar(s,x,COL_Y,cw,mod.get('name',''),DARK,hh,12)
        status=mod.get('status','')
        col,tcol=STATUS_CHIP.get(status,(GRAY,WHITE))
        cw2,chh=1.3,0.3; cx=x+cw-cw2-0.14; cyc=COL_Y+(hh-chh)/2
        chip=rrect(s,cx,cyc,cw2,chh,col,radius=0.5)
        ctf=chip.text_frame;ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
        cpp=ctf.paragraphs[0];crr=cpp.add_run();crr.text=status;crr.font.size=Pt(9);crr.font.bold=True;crr.font.color.rgb=tcol;crr.font.name="Arial"
        tf=tbox(s,x+0.24,COL_Y+hh+0.2,cw-0.48,COL_H-hh-0.4)
        for j,(ml,mv) in enumerate(mod.get('metrics') or []):
            line(tf,[(mv,20,True,WHITE)],first=(j==0),sb=(0 if j==0 else 12))
            line(tf,[(ml,10,False,SUB)],sb=2)
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Source: TrendAI Vision One — Zero Trust Secure Access, Secure Access Overview.",8,False,SUB)],first=True)

# ========================= 7. CREDITS =========================
def build_credits(s,D):
    d=D['credits']
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("Credits",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   TrendAI™ Flex Licensing · Platform Usage and Credits   •   {d['term']}   •   as of {D['as_of']}",11,False,SUB)],first=True)
    rrect(s,0.5,1.4,12.33,1.05,CARD,BORDER,radius=0.05)
    def hstat(cx,num,lab,col):
        tf=tbox(s,cx-1.95,1.43,3.9,0.99); line(tf,[(num,26,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    hstat(2.04,d['purchased'],"Credits purchased",WHITE)
    hstat(5.15,d['used'],f"Used to date  ({d['used_pct']})",REDB)
    hstat(8.26,d['balance'],f"Balance remaining  ({d['balance_pct']})",GREEN)
    hstat(11.37,d['planned'],"Planned credits / month",BLUE)
    for vx in [3.59,6.70,9.81]: rrect(s,vx,1.56,0.014,0.73,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    px,pw=0.5,7.4; py,ph=2.6,4.55
    rrect(s,px,py,pw,ph,CARD,BORDER,radius=0.04)
    headerbar(s,px,py,pw,"PLANNED USAGE vs ACTUAL DRAWDOWN  ·  LAST 3 MONTHS",RED,0.5,11.5)
    def legdot(x,y,c,label): rrect(s,x,y+0.06,0.16,0.16,c,radius=0.5); line(tbox(s,x+0.2,y,1.3,0.28),[(label,9.5,False,TXT)],first=True)
    legdot(px+0.7,py+0.56,BLUE,"Planned"); legdot(px+2.1,py+0.56,REDB,"Actual")
    data=d['months']  # [month,planned,actual]
    maxY=float(d['chart_max'])
    plotL=px+0.72; plotR=px+pw-0.25; plotW=plotR-plotL
    baseY=py+ph-0.55; topY=py+1.05; plotH=baseY-topY; scale=plotH/maxY
    for g in d['gridlines']:
        gy=baseY-g*scale
        rrect(s,plotL-0.05,gy,plotW+0.1,0.01,GRID,radius=0,shape=MSO_SHAPE.RECTANGLE)
        line(tbox(s,px+0.1,gy-0.12,0.55,0.24),[((str(g//1000)+"k" if g else "0"),8.5,False,SUB)],align=PP_ALIGN.RIGHT,first=True)
    gW=plotW/len(data); colW=0.66; gap=0.16
    for i,(mon,pl,ac) in enumerate(data):
        gx=plotL+i*gW; cx1=gx+gW/2-colW-gap/2; cx2=gx+gW/2+gap/2
        h1=pl*scale; rrect(s,cx1,baseY-h1,colW,h1,BLUE,radius=0.12)
        line(tbox(s,cx1-0.1,baseY-h1-0.26,colW+0.2,0.24),[(f"{pl:,}",9,True,WHITE)],align=PP_ALIGN.CENTER,first=True)
        h2=ac*scale; rrect(s,cx2,baseY-h2,colW,h2,REDB,radius=0.12)
        line(tbox(s,cx2-0.1,baseY-h2-0.26,colW+0.2,0.24),[(f"{ac:,}",9,True,WHITE)],align=PP_ALIGN.CENTER,first=True)
        line(tbox(s,gx,baseY+0.06,gW,0.28),[(mon,10,True,TXT)],align=PP_ALIGN.CENTER,first=True)
    ex,ew=8.1,4.73
    rrect(s,ex,py,ew,ph,CARD,BORDER,radius=0.04)
    headerbar(s,ex,py,ew,"PURCHASED vs USED",BLUE,0.5,11.5)
    etf=tbox(s,ex+0.22,py+0.62,ew-0.44,ph-0.78,MSO_ANCHOR.TOP)
    for i,b in enumerate(d['bullets']):
        line(etf,[("•  ",10.5,True,BLUE),(b,10.5,False,TXT)],first=(i==0),sb=(0 if i==0 else 8))
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Planned usage = purchased monthly plan (baseline); actual drawdown = credits deducted each month. Source: TrendAI Flex Licensing — Platform Usage and Credits.",8,False,SUB)],first=True)

# ========================= 8. END OF YEAR PREDICTION =========================
def nice_ceiling(v):
    """Smallest of {1,2,2.5,5,10} x 10^n that is >= v (n picked as needed)."""
    if v<=0: return 100
    magnitude=10**math.floor(math.log10(v))
    for mult in (1,2,2.5,5,10):
        cand=mult*magnitude
        if cand>=v: return cand
    return 10*magnitude

def build_eoy_prediction(s,D):
    d=D['credits']
    num=lambda x:int(str(x).replace(',',''))
    purchased=num(d['purchased'])
    actuals=[m[2] for m in d['months']]               # actual monthly drawdowns
    monthly=round(sum(actuals)/len(actuals))           # avg monthly run-rate
    predicted=monthly*12                               # "the monthly times 12"
    surplus=purchased-predicted
    surplus_pct=round(surplus/purchased*100)
    used_pct=round(predicted/purchased*100)
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.28,12.33,0.58,MSO_ANCHOR.TOP),[("End of Year Prediction Summary",30,True,WHITE)],first=True)
    rrect(s,0.55,0.9,3.4,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.97,12.33,0.3,MSO_ANCHOR.TOP),[(f"{CUST}   •   Full-year credit projection (monthly run-rate × 12)   •   {d['term']}   •   as of {D['as_of']}",11,False,SUB)],first=True)
    rrect(s,0.5,1.4,12.33,1.05,CARD,BORDER,radius=0.05)
    def hstat(cx,n,lab,col):
        tf=tbox(s,cx-1.95,1.43,3.9,0.99); line(tf,[(n,26,True,col)],align=PP_ALIGN.CENTER,first=True); line(tf,[(lab,10.5,False,SUB)],align=PP_ALIGN.CENTER,sb=2)
    hstat(2.55,f"{purchased:,}","Credits purchased (annual)",WHITE)
    hstat(6.665,f"{predicted:,}","Predicted used (full year)",REDB)
    hstat(10.78,f"{surplus:,}",f"Projected surplus  ({surplus_pct}%)",GREEN)
    for vx in [4.61,8.72]: rrect(s,vx,1.56,0.014,0.73,BORDER,radius=0,shape=MSO_SHAPE.RECTANGLE)
    # ---- bar chart: Purchased vs Predicted Used ----
    px,pw=0.5,7.4; py,ph=2.6,4.55
    rrect(s,px,py,pw,ph,CARD,BORDER,radius=0.04)
    headerbar(s,px,py,pw,"PURCHASED vs PREDICTED ANNUAL USAGE",RED,0.5,11.5)
    maxY=nice_ceiling(max(purchased,predicted)*1.05); grid=[maxY*i//4 for i in range(5)]
    plotL=px+0.95; plotR=px+pw-0.35; plotW=plotR-plotL
    baseY=py+ph-0.6; topY=py+1.0; plotH=baseY-topY; scale=plotH/maxY
    def fmt_tick(g):
        if not g: return "0"
        if g>=1000000: return f"{g/1000000:g}M"
        return f"{g//1000}k"
    for g in grid:
        gy=baseY-g*scale
        rrect(s,plotL-0.05,gy,plotW+0.1,0.01,GRID,radius=0,shape=MSO_SHAPE.RECTANGLE)
        line(tbox(s,px+0.12,gy-0.12,0.72,0.24),[(fmt_tick(g),8.5,False,SUB)],align=PP_ALIGN.RIGHT,first=True)
    bars=[("Purchased (annual)",purchased,BLUE),("Predicted Used (FY)",predicted,REDB)]
    colW=1.5; gW=plotW/len(bars)
    for i,(lab,val,col) in enumerate(bars):
        cx=plotL+i*gW+gW/2-colW/2; h=val*scale
        rrect(s,cx,baseY-h,colW,h,col,radius=0.12)
        line(tbox(s,cx-0.3,baseY-h-0.32,colW+0.6,0.28),[(f"{val:,}",13.5,True,WHITE)],align=PP_ALIGN.CENTER,first=True)
        line(tbox(s,plotL+i*gW,baseY+0.08,gW,0.3),[(lab,10.5,True,TXT)],align=PP_ALIGN.CENTER,first=True)
    # ---- year-end position panel ----
    ex,ew=8.1,4.73
    rrect(s,ex,py,ew,ph,CARD,BORDER,radius=0.04)
    headerbar(s,ex,py,ew,"PROJECTED YEAR-END POSITION",BLUE,0.5,11.5)
    stf=tbox(s,ex+0.25,py+0.72,ew-0.5,1.05)
    line(stf,[(f"{surplus:,}",32,True,GREEN)],first=True)
    line(stf,[(f"credits projected unused at term end  (~{surplus_pct}%)",10.5,False,SUB)],sb=2)
    btf=tbox(s,ex+0.25,py+2.15,ew-0.5,ph-2.35,MSO_ANCHOR.TOP)
    bullets=[
        f"Annual purchase: {purchased:,} credits for the {d['term'].replace('Term: ','')} term.",
        f"Average monthly drawdown: ~{monthly:,} credits ({' → '.join(f'{a:,}' for a in actuals)} over the last {len(actuals)} months).",
        f"Run-rate projection: {monthly:,} × 12 = {predicted:,} credits used by year end.",
        f"At this pace {CUST} consumes ~{used_pct}% of purchased credits, leaving a ~{surplus:,}-credit buffer.",
    ]
    for i,b in enumerate(bullets):
        line(btf,[("•  ",10.5,True,BLUE),(b,10.5,False,TXT)],first=(i==0),sb=(0 if i==0 else 9))
    line(tbox(s,0.5,7.22,12.33,0.24,MSO_ANCHOR.TOP),[("Prediction = average actual monthly drawdown × 12 months. Source: TrendAI Flex Licensing — Platform Usage and Credits.",8,False,SUB)],first=True)

# ========================= 9. ACCOUNT TEAM =========================
def build_team(s,D):
    rrect(s,0,0,13.333,7.5,BG,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,0.30,12.33,0.6,MSO_ANCHOR.TOP),[("TrendAI Account Team",30,True,WHITE)],first=True)
    rrect(s,0.55,0.97,2.6,0.05,RED,radius=0,shape=MSO_SHAPE.RECTANGLE)
    line(tbox(s,0.5,1.04,12.33,0.32,MSO_ANCHOR.TOP),[(f"{CUST}   •   Your dedicated Trend Micro team",11,False,SUB)],first=True)
    colmap={"red":RED,"teal":TEAL,"blue":BLUE,"orange":ORANGE2}
    team=D['team']['roles']  # [role,colorkey]
    cw=5.96; chh=2.30; gx=0.41; gy=0.34; xs=[0.5,0.5+cw+gx]; ys=[1.55,1.55+chh+gy]
    pos=[(xs[0],ys[0]),(xs[1],ys[0]),(xs[0],ys[1]),(xs[1],ys[1])]
    for (role,ck),(cx,cy) in zip(team,pos):
        accent=colmap[ck]
        rrect(s,cx,cy,cw,chh,CARD,BORDER,radius=0.05)
        rrect(s,cx,cy,0.08,chh,accent,radius=0,shape=MSO_SHAPE.RECTANGLE)
        av=cy+(chh-1.1)/2
        oval(s,cx+0.32,av,1.1,1.1,CARD2,line=accent,lw=2.0)
        line(tbox(s,cx+0.32,av,1.1,1.1),[("…",26,True,accent)],align=PP_ALIGN.CENTER,first=True)
        tf=tbox(s,cx+1.66,cy+0.2,cw-1.9,chh-0.4)
        line(tf,[(role,10,True,accent)],first=True)
        line(tf,[("[ Full Name ]",17,True,WHITE)],sb=3)
        line(tf,[("Email   ",10.5,False,SUB),("[ name@trendmicro.com ]",10.5,False,PH)],sb=7)
        line(tf,[("Phone   ",10.5,False,SUB),("[ +1 (___) ___-____ ]",10.5,False,PH)],sb=4)
    fy=6.62
    rrect(s,0.5,fy,12.33,0.6,CARD2,BORDER,radius=0.08)
    line(tbox(s,0.78,fy,11.8,0.6),[("Reach out to your account team any time with questions about this report, your environment, or your TrendAI Vision One™ subscription.",10.5,False,TXT)],first=True)

# ========================= MAIN =========================
def main():
    global CUST
    with open(DATA_FILE) as f: D=json.load(f)
    CUST=D['customer']
    OUTPUT = os.path.expanduser(f"~/Desktop/{CUST} Value Report.pptx")
    # Content slides only (slide 1 title + last slide Account Team are managed separately).
    # Risk Factors Detail slides are generated dynamically, 2 categories per slide,
    # from D['risk_detail'] (a flat list — length can vary run to run).
    risk_detail=D.get('risk_detail',[])
    pairs=[risk_detail[i:i+2] for i in range(0,len(risk_detail),2)]
    risk_detail_builders=[(lambda s,D,pair=pair: build_risk_detail_slide(s,D,pair)) for pair in pairs]
    content_builders=[build_cyber]+risk_detail_builders+[build_mdr,build_data_source,build_intel,build_workflow_automation,build_ai_security,build_zero_trust,build_identity,build_data_security,build_standard_endpoint_protection,build_server_workload_protection,build_cloud_security,build_network_security,build_email,build_credits,build_eoy_prediction]
    expected=1+len(content_builders)+1   # title + content + Account Team

    # PRESERVE slide 1 (title) and the LAST slide (Account Team): if a prior deck
    # already exists with the expected layout, reopen it and rebuild ONLY the
    # content slides in place, so any hand-edits (presenter name, team contacts)
    # to those two slides survive reruns.
    if os.path.exists(OUTPUT):
        prs=Presentation(OUTPUT)
        slides=list(prs.slides)
        if len(slides)==expected:
            for i,b in enumerate(content_builders):
                sl=slides[i+1]          # slides[0]=title and slides[-1]=Account Team left intact
                clear(sl); b(sl,D)
            update_title_date(slides[0],D)   # refresh ONLY the date on the preserved title slide
            prs.save(OUTPUT)
            print("Wrote",OUTPUT,"- rebuilt",len(content_builders),
                  "content slides; preserved slide 1 + last slide (date refreshed)")
            return
        print("Existing deck has",len(slides),"slides (expected",expected,
              ")- doing a full rebuild from template")

    # First run / unexpected layout: full build from template.
    prs=Presentation(TEMPLATE)
    quote=None
    for m in prs.slide_masters:
        for l in m.slide_layouts:
            if l.name=="Quote 01": quote=l
    if quote is None: raise SystemExit("Quote 01 layout not found in template")
    build_title(prs.slides[0],D)        # title slide (already in template as slide 0)
    for b in content_builders+[build_team]:
        sl=prs.slides.add_slide(quote)
        clear(sl)
        b(sl,D)
    prs.save(OUTPUT)
    print("Wrote",OUTPUT,"(full build) with",len(list(prs.slides)),"slides")

if __name__=="__main__":
    main()
